"""
修改ppt每一页第一个文本框的颜色
"""
import pptx
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
from pptx.util import Cm, Pt

def alter_color(ppt):


    num = len(ppt.slides)
    for i in range(1,num):
        slide = ppt.slides[i]
        shape = slide.shapes[0]

        print(shape.text)
        shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 0)


def alter_bold(ppt):

    num = len(ppt.slides)
    for i in range(1, num):
        slide = ppt.slides[i]
        shape = slide.shapes[-1]

        print(shape.text)
        shape.text_frame.paragraphs[0].runs[0].font.bold = True

    ppt.save("六级词汇.pptx")

def add_textbox_me(slide,left,top,width,height,text_content,font_size = 24,font_name = "微软雅黑 (标题)",text_color = (255, 255, 255)):
    """测试写入文本框  index指的是第几个文本框，根据文本框序号确定字体、字号以及颜色"""
    text_box_1 = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box_1.text_frame
    tf.word_wrap = True
    tf.text = text_content
    tf.paragraphs[0].font.bold = True
    # print(tf.text)
    # slide.shapes[index].text_frame.paragraphs[0].font.name = "微软雅黑 (标题)"
    # slide.shapes[index].text_frame.paragraphs[0].font.size = Pt(24)
    # if (index == 0):
    #     slide.shapes[index].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
    # else:
    #     slide.shapes[index].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

def add_textbox_cet6_words_yaru(ppt):

    num = len(ppt.slides)

    wb = load_workbook(r"D:\工作\神龙录题 强基+四六级听力\神龙录题 强基+四六级听力\2. 六级ppt\六级词汇制作ppt用 .xlsx")
    sheet = wb["制作ppt用总表-雅茹"]
    left = Cm(3.33)
    top = Cm(13.73)
    width = Cm(18.4)
    height = Cm(2.31)
    for i in range(2,22):
        content = sheet.cell(i,13).value
        # list_contents.append(content)
        print(content)
        slide = ppt.slides[i-1]
        add_textbox_me(slide, left, top, width, height, content)

def main():
    ppt = pptx.Presentation(pptx="六级词汇-雅茹 - 副本.pptx")
    alter_color(ppt)
    ppt.save("六级词汇-雅茹 - 副本.pptx")

if __name__ == '__main__':
    main()